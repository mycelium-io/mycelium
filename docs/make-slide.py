#!/usr/bin/env python3
"""
Generate the dissent-map opening-positions slide as a .pptx file.

Usage:
  pip install python-pptx
  python make-slide.py --out docs/dissent-positions.pptx
"""

import argparse
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Palette (matches Mycelium dark UI) ─────────────────────────────────────────

BG          = RGBColor(0x0D, 0x0D, 0x0D)
ACCENT      = RGBColor(0x5D, 0xD4, 0xE0)   # teal
MUTED       = RGBColor(0x55, 0x55, 0x66)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DIM         = RGBColor(0xAA, 0xAA, 0xBB)
BORDER      = RGBColor(0x22, 0x22, 0x33)

PM_COLOR    = RGBColor(0xFB, 0xBF, 0x24)   # amber  — business urgency
ENG_COLOR   = RGBColor(0x34, 0xD3, 0x99)   # green  — technical correctness
SRE_COLOR   = RGBColor(0x60, 0xA5, 0xFA)   # blue   — reliability
SEC_COLOR   = RGBColor(0xF8, 0x71, 0x71)   # red    — risk / security

# ── Slide dimensions (1920×1080 px at 96 dpi ≈ 20×11.25 in) ───────────────────

W = Inches(20)
H = Inches(11.25)

# ── Agent definitions ─────────────────────────────────────────────────────────

AGENTS = [
    {
        "handle": "pm-agent",
        "role":   "Product Manager",
        "color":  PM_COLOR,
        "pos":    "top-left",
        "constraint": (
            "Friday launch is committed.\n"
            "Slipping cancels a paid marketing campaign\n"
            "and triggers the enterprise SLA penalty."
        ),
        "stance": "SHIP FRIDAY",
    },
    {
        "handle": "eng-lead",
        "role":   "Engineering Lead",
        "color":  ENG_COLOR,
        "pos":    "top-right",
        "constraint": (
            "Will not ship known data corruption.\n"
            "Bug corrupts order metadata in ~3%\n"
            "of discount-code orders."
        ),
        "stance": "FIX FIRST",
    },
    {
        "handle": "sre-agent",
        "role":   "Site Reliability",
        "color":  SRE_COLOR,
        "pos":    "bottom-left",
        "constraint": (
            "Ship only with a kill switch on checkout.\n"
            "Must be reversible in under 2 minutes\n"
            "or bad orders compound."
        ),
        "stance": "KILL SWITCH",
    },
    {
        "handle": "secops-agent",
        "role":   "Security Ops",
        "color":  SEC_COLOR,
        "pos":    "bottom-right",
        "constraint": (
            "Cannot sign off until exploitability\n"
            "is confirmed. Discount-code bugs are\n"
            "a known attack surface."
        ),
        "stance": "REVIEW FIRST",
    },
]

# Corner positions (left, top) in inches for a 20×11.25 slide
CORNER_POSITIONS = {
    "top-left":     (0.4,  0.5),
    "top-right":    (13.8, 0.5),
    "bottom-left":  (0.4,  7.2),
    "bottom-right": (13.8, 7.2),
}
BOX_W = Inches(5.8)
BOX_H = Inches(3.3)


def set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size, bold=False,
                color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def add_agent_box(slide, agent: dict):
    lft_in, top_in = CORNER_POSITIONS[agent["pos"]]
    left = Inches(lft_in)
    top  = Inches(top_in)
    color = agent["color"]

    # Outer border box
    border = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, BOX_W, BOX_H,
    )
    border.fill.solid()
    border.fill.fore_color.rgb = RGBColor(0x12, 0x12, 0x1E)
    border.line.color.rgb = color
    border.line.width = Pt(1.5)

    pad = Inches(0.2)

    # Stance badge (top strip)
    badge = slide.shapes.add_shape(
        1,
        left, top, BOX_W, Inches(0.45),
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    tf = badge.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = agent["stance"]
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x0D, 0x0D, 0x0D)

    # Handle
    add_textbox(
        slide,
        left + pad, top + Inches(0.55),
        BOX_W - pad * 2, Inches(0.4),
        f"@{agent['handle']}",
        font_size=13, bold=True, color=color,
    )

    # Role
    add_textbox(
        slide,
        left + pad, top + Inches(0.95),
        BOX_W - pad * 2, Inches(0.3),
        agent["role"],
        font_size=10, color=DIM,
    )

    # Constraint text
    add_textbox(
        slide,
        left + pad, top + Inches(1.35),
        BOX_W - pad * 2, Inches(1.8),
        agent["constraint"],
        font_size=11, color=WHITE, wrap=True,
    )


def add_center_topic(slide):
    cx = Inches(6.3)
    cy = Inches(3.6)
    cw = Inches(7.4)
    ch = Inches(4.0)

    # Center box
    box = slide.shapes.add_shape(1, cx, cy, cw, ch)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x10, 0x10, 0x1C)
    box.line.color.rgb = ACCENT
    box.line.width = Pt(1)

    # Title
    add_textbox(
        slide, cx, cy + Inches(0.3), cw, Inches(0.6),
        "FRIDAY LAUNCH — SHIP OR SLIP?",
        font_size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER,
    )

    # Divider line (thin rectangle)
    div = slide.shapes.add_shape(1,
        cx + Inches(0.5), cy + Inches(1.05),
        cw - Inches(1.0), Inches(0.02),
    )
    div.fill.solid()
    div.fill.fore_color.rgb = MUTED
    div.line.fill.background()

    # Scenario text
    scenario = (
        "A checkout bug is found 24 hours before a major Friday launch.\n\n"
        "Four autonomous agents. Four hard constraints.\n"
        "No central authority. One decision required."
    )
    add_textbox(
        slide, cx + Inches(0.4), cy + Inches(1.2),
        cw - Inches(0.8), Inches(2.0),
        scenario,
        font_size=12, color=DIM, align=PP_ALIGN.CENTER, wrap=True,
    )

    # Footer label
    add_textbox(
        slide, cx, cy + Inches(3.3), cw, Inches(0.4),
        "mycelium  ·  IoC dissent-map demo",
        font_size=9, color=MUTED, align=PP_ALIGN.CENTER,
    )


def build(out_path: str):
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    blank_layout = prs.slide_layouts[6]  # completely blank
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, BG)

    # Top label
    add_textbox(
        slide,
        Inches(6.3), Inches(0.15), Inches(7.4), Inches(0.4),
        "OPENING POSITIONS",
        font_size=10, color=MUTED, align=PP_ALIGN.CENTER,
    )

    for agent in AGENTS:
        add_agent_box(slide, agent)

    add_center_topic(slide)

    prs.save(out_path)
    print(f"✅ Saved: {out_path}")


HTML_TEMPLATE = """<!DOCTYPE html>
<html>
<head>
<meta charset="utf-8">
<style>
  * { margin: 0; padding: 0; box-sizing: border-box; }
  body {
    width: 1920px; height: 1080px; overflow: hidden;
    background: #0d0d0d;
    font-family: 'SF Mono', 'Fira Code', 'Consolas', monospace;
    color: #fff;
  }
  .slide { position: relative; width: 1920px; height: 1080px; }

  .label-top {
    position: absolute; top: 14px; left: 0; right: 0;
    text-align: center; font-size: 13px; color: #555566;
    letter-spacing: 0.15em; text-transform: uppercase;
  }

  /* Agent boxes */
  .agent {
    position: absolute; width: 560px; height: 318px;
    background: #12121e;
    border: 1.5px solid var(--color);
    display: flex; flex-direction: column;
  }
  .agent .badge {
    background: var(--color); color: #0d0d0d;
    font-size: 12px; font-weight: 700; letter-spacing: 0.12em;
    text-align: center; padding: 8px 0; text-transform: uppercase;
  }
  .agent .body { padding: 14px 18px; }
  .agent .handle { font-size: 14px; font-weight: 700; color: var(--color); margin-bottom: 3px; }
  .agent .role { font-size: 11px; color: #aaaabb; margin-bottom: 12px; }
  .agent .constraint { font-size: 12px; color: #ddddee; line-height: 1.6; }

  /* Positions */
  .top-left     { top: 48px;  left: 38px;  }
  .top-right    { top: 48px;  right: 38px; }
  .bottom-left  { bottom: 48px; left: 38px;  }
  .bottom-right { bottom: 48px; right: 38px; }

  /* Center box */
  .center {
    position: absolute;
    left: 606px; top: 346px; width: 708px; height: 386px;
    background: #10101c;
    border: 1px solid #5dd4e0;
    display: flex; flex-direction: column; align-items: center;
    padding: 28px 36px;
  }
  .center .title {
    font-size: 17px; font-weight: 700; color: #5dd4e0;
    letter-spacing: 0.08em; text-align: center; margin-bottom: 16px;
  }
  .center .divider {
    width: 80%; height: 1px; background: #333344; margin-bottom: 20px;
  }
  .center .scenario {
    font-size: 13px; color: #aaaabb; line-height: 1.8;
    text-align: center;
  }
  .center .footer {
    position: absolute; bottom: 16px;
    font-size: 10px; color: #333344; letter-spacing: 0.1em;
  }

  /* Connector lines (SVG overlay) */
  svg.connectors {
    position: absolute; top: 0; left: 0;
    width: 1920px; height: 1080px;
    pointer-events: none;
  }
</style>
</head>
<body>
<div class="slide">

  <div class="label-top">OPENING POSITIONS</div>

  <!-- SVG connector lines from corners to center -->
  <svg class="connectors">
    <line x1="598"  y1="207" x2="606"  y2="539" stroke="#222233" stroke-width="1" stroke-dasharray="6,4"/>
    <line x1="1322" y1="207" x2="1314" y2="539" stroke="#222233" stroke-width="1" stroke-dasharray="6,4"/>
    <line x1="598"  y1="874" x2="606"  y2="732" stroke="#222233" stroke-width="1" stroke-dasharray="6,4"/>
    <line x1="1322" y1="874" x2="1314" y2="732" stroke="#222233" stroke-width="1" stroke-dasharray="6,4"/>
  </svg>

  <!-- pm-agent: top-left -->
  <div class="agent top-left" style="--color: #fbbf24">
    <div class="badge">SHIP FRIDAY</div>
    <div class="body">
      <div class="handle">@pm-agent</div>
      <div class="role">Product Manager</div>
      <div class="constraint">Friday launch is committed.<br>
        Slipping cancels a paid marketing campaign<br>
        and triggers the enterprise SLA penalty.</div>
    </div>
  </div>

  <!-- eng-lead: top-right -->
  <div class="agent top-right" style="--color: #34d399">
    <div class="badge">FIX FIRST</div>
    <div class="body">
      <div class="handle">@eng-lead</div>
      <div class="role">Engineering Lead</div>
      <div class="constraint">Will not ship known data corruption.<br>
        Bug corrupts order metadata in ~3%<br>
        of discount-code orders.</div>
    </div>
  </div>

  <!-- sre-agent: bottom-left -->
  <div class="agent bottom-left" style="--color: #60a5fa">
    <div class="badge">KILL SWITCH</div>
    <div class="body">
      <div class="handle">@sre-agent</div>
      <div class="role">Site Reliability</div>
      <div class="constraint">Ship only with a kill switch on checkout.<br>
        Must be reversible in under 2 minutes<br>
        or bad orders will compound.</div>
    </div>
  </div>

  <!-- secops-agent: bottom-right -->
  <div class="agent bottom-right" style="--color: #f87171">
    <div class="badge">REVIEW FIRST</div>
    <div class="body">
      <div class="handle">@secops-agent</div>
      <div class="role">Security Ops</div>
      <div class="constraint">Cannot sign off until exploitability<br>
        is confirmed. Discount-code bugs are<br>
        a known attack surface.</div>
    </div>
  </div>

  <!-- Center -->
  <div class="center">
    <div class="title">FRIDAY LAUNCH — SHIP OR SLIP?</div>
    <div class="divider"></div>
    <div class="scenario">
      A checkout bug is found 24 hours before a major Friday launch.<br><br>
      Four autonomous agents. Four hard constraints.<br>
      No central authority. One decision required.
    </div>
    <div class="footer">mycelium &nbsp;·&nbsp; IoC dissent-map demo</div>
  </div>

</div>
</body>
</html>
"""


def build_html(out_path: str):
    with open(out_path, "w") as f:
        f.write(HTML_TEMPLATE)
    print(f"✅ HTML saved: {out_path}")


if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("--out", default="docs/dissent-positions.pptx")
    parser.add_argument("--html", default="docs/dissent-positions.html")
    args = parser.parse_args()
    build(args.out)
    build_html(args.html)
